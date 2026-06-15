"""Generate nature-paper2ppt revised Chinese PPTX from Agentic AI literature review."""

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from create_ai_paper import (
    CHARTS_DIR,
    chart_agent_loop,
    chart_agentic_taxonomy,
    chart_framework_comparison,
    chart_timeline,
)

PROJECT_DIR = Path(__file__).parent
OUTPUT = PROJECT_DIR / "agentic_ai_presentation.pptx"
ASSETS_DIR = PROJECT_DIR / "output" / "assets" / "figures"
QA_REPORT = PROJECT_DIR / "output" / "qa_report.md"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x1E, 0x3A, 0x5F)
BLUE = RGBColor(0x25, 0x63, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x64, 0x74, 0x8B)
ACCENT = RGBColor(0x38, 0xBD, 0xF8)


def set_slide_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_speaker_notes(slide, text: str):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def add_footer(slide, text: str, prs: Presentation):
    box = slide.shapes.add_textbox(
        Inches(0.5), prs.slide_height - Inches(0.45),
        prs.slide_width - Inches(1.0), Inches(0.35),
    )
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def add_title_bar(slide, title: str, prs: Presentation, size=26):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.18), prs.slide_width - Inches(1.1), Inches(0.85))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = WHITE


def add_bullets(slide, items, left=Inches(0.65), top=Inches(1.35), width=Inches(12.1),
                height=Inches(5.2), size=16):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)


def add_takeaway_strip(slide, text: str, prs: Presentation):
    box = slide.shapes.add_textbox(Inches(0.55), prs.slide_height - Inches(1.35), Inches(12.2), Inches(0.45))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = BLUE


def add_two_column_bullets(slide, left_title, left_items, right_title, right_items):
    lt = slide.shapes.add_textbox(Inches(0.55), Inches(1.25), Inches(5.9), Inches(0.45))
    lt.text_frame.text = left_title
    lt.text_frame.paragraphs[0].font.size = Pt(15)
    lt.text_frame.paragraphs[0].font.bold = True
    lt.text_frame.paragraphs[0].font.color.rgb = BLUE
    add_bullets(slide, left_items, left=Inches(0.55), top=Inches(1.7), width=Inches(5.9), height=Inches(4.5), size=15)

    rt = slide.shapes.add_textbox(Inches(6.85), Inches(1.25), Inches(5.9), Inches(0.45))
    rt.text_frame.text = right_title
    rt.text_frame.paragraphs[0].font.size = Pt(15)
    rt.text_frame.paragraphs[0].font.bold = True
    rt.text_frame.paragraphs[0].font.color.rgb = BLUE
    add_bullets(slide, right_items, left=Inches(6.85), top=Inches(1.7), width=Inches(5.9), height=Inches(4.5), size=15)


def add_table_slide(slide, headers, rows, top=Inches(1.45), font_size=12):
    cols = len(headers)
    row_h = Inches(0.42)
    tbl_shape = slide.shapes.add_table(len(rows) + 1, cols, Inches(0.55), top, Inches(12.2), row_h)
    table = tbl_shape.table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(font_size)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = DARK
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
    for row in table.rows:
        row.height = row_h


def add_figure_slide(slide, image_path: Path, caption: str, source: str, prs: Presentation,
                     img_width=Inches(9.0), top=Inches(1.35)):
    slide.shapes.add_picture(str(image_path), (prs.slide_width - img_width) / 2, top, width=img_width)
    cap = slide.shapes.add_textbox(Inches(0.55), prs.slide_height - Inches(0.95), Inches(12.2), Inches(0.35))
    cap.text_frame.text = caption
    cap.text_frame.paragraphs[0].font.size = Pt(11)
    cap.text_frame.paragraphs[0].font.italic = True
    cap.text_frame.paragraphs[0].font.color.rgb = MUTED
    cap.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    src = slide.shapes.add_textbox(Inches(0.55), prs.slide_height - Inches(0.55), Inches(12.2), Inches(0.25))
    src.text_frame.text = source
    src.text_frame.paragraphs[0].font.size = Pt(8)
    src.text_frame.paragraphs[0].font.color.rgb = MUTED
    src.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


def stage_assets(chart_paths: dict[str, Path]) -> dict[str, Path]:
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)
    staged = {}
    for name, src in chart_paths.items():
        dest = ASSETS_DIR / src.name
        shutil.copy2(src, dest)
        staged[name] = dest
    return staged


def audit_presentation(path: Path) -> list[str]:
    prs = Presentation(str(path))
    issues = []
    sw, sh = prs.slide_width, prs.slide_height
    notes_count = 0
    media_count = 0

    for i, slide in enumerate(prs.slides, start=1):
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            notes_count += 1
        char_total = 0
        text_boxes = 0
        for shape in slide.shapes:
            if shape.shape_type == 13:  # picture
                media_count += 1
            if not shape.has_text_frame:
                continue
            text_boxes += 1
            text = shape.text_frame.text or ""
            char_total += len(text)
            if len(text) > 120 and shape.top < Inches(4):
                issues.append(f"medium: slide {i} text box may be dense ({len(text)} chars)")
            l, t, w, h = shape.left, shape.top, shape.width, shape.height
            if l < 0 or t < 0 or l + w > sw + Emu(91440) or t + h > sh + Emu(91440):
                issues.append(f"high: slide {i} shape out of bounds")
        if char_total > 280:
            issues.append(f"medium: slide {i} total on-slide text high ({char_total} chars)")

    if notes_count < len(prs.slides) - 1:
        issues.append(f"medium: speaker notes on {notes_count}/{len(prs.slides)} slides")
    return [
        f"slides: {len(prs.slides)}",
        f"embedded media: {media_count}",
        f"slides with notes: {notes_count}",
        *issues,
    ]


def write_qa_report(audit_lines: list[str], total: int):
    QA_REPORT.parent.mkdir(parents=True, exist_ok=True)
    fixed = [
        "Converted deck to simplified Chinese with conclusion-style titles",
        "Restructured narrative to review/evidence-map arc (nature-paper2ppt)",
        "Shortened on-slide bullets; moved detail to speaker notes",
        "Added source labels on figure slides",
        "Split dense challenge content across two slides",
    ]
    body = f"""# QA Report — agentic_ai_presentation.pptx

## Status
- PPTX created: yes
- Skill: nature-paper2ppt (paper_type: review / evidence-map)
- Slide count: {total}
- Language: simplified Chinese (technical terms preserved in English)

## Verification
{chr(10).join('- ' + line for line in audit_lines)}

## Self-review corrections applied
{chr(10).join('- ' + line for line in fixed)}

## Known limitations
- Charts are redrawn schematics (改绘自 literature review figures), not paper originals
- Framework maturity chart based on MASEval (2026) synthesis; verify before citation
"""
    QA_REPORT.write_text(body, encoding="utf-8")


def build_presentation():
    CHARTS_DIR.mkdir(exist_ok=True)
    charts = {
        "taxonomy": chart_agentic_taxonomy(),
        "loop": chart_agent_loop(),
        "timeline": chart_timeline(),
        "framework": chart_framework_comparison(),
    }
    assets = stage_assets(charts)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    total = 15

    def footer(n):
        return f"{n} / {total}"

    # --- 1 标题页 ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, NAVY)
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.15), prs.slide_width, Inches(0.07))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    title = s.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(11.9), Inches(1.6))
    title.text_frame.text = "智能体人工智能"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = WHITE
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    sub = s.shapes.add_textbox(Inches(0.7), Inches(2.35), Inches(11.9), Inches(0.65))
    sub.text_frame.text = "Agentic AI · 文献综述"
    sub.text_frame.paragraphs[0].font.size = Pt(24)
    sub.text_frame.paragraphs[0].font.color.rgb = ACCENT
    sub.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    meta = s.shapes.add_textbox(Inches(0.7), Inches(3.55), Inches(11.9), Inches(1.0))
    meta.text_frame.text = "计算机科学与工程 · 2026年6月"
    meta.text_frame.paragraphs[0].font.size = Pt(16)
    meta.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    meta.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_speaker_notes(s, "开场：本汇报基于 Agentic AI 文献综述，采用 evidence-map 叙事结构，约15分钟。")

    # --- 2 为什么现在重要 ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_title_bar(s, "LLM 已能推理，但尚难自主完成多步目标", prs)
    add_bullets(s, [
        "生成式 AI：单次 prompt → 单次输出",
        "智能体 AI：感知环境、规划、调用工具、迭代至目标",
        "2022–2025 年 ReAct、MCP、多智能体框架快速涌现",
        "产业落地需求推动从「会答」到「会做」的范式转变",
    ])
    add_takeaway_strip(s, "核心转变：从文本生成到闭环控制", prs)
    add_footer(s, footer(2), prs)
    add_speaker_notes(s, "强调时机：LLM 能力已具备，但缺乏自主多步执行与工具编排；智能体是连接模型与真实任务的桥梁。")

    # --- 3 概念框架 ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_title_bar(s, "智能体 AI 由推理·行动·交互三轴能力构成", prs)
    add_figure_slide(
        s, assets["taxonomy"],
        "推理（规划/记忆）、行动（工具/编排）、交互（多智能体/人机协同）",
        "改绘自 Xi et al., 2025 能力分类",
        prs, img_width=Inches(10.5), top=Inches(1.25),
    )
    add_footer(s, footer(3), prs)
    add_speaker_notes(s, "用 taxonomy 图作为全场概念地图；三轴不是并列模块，而是智能体闭环的不同切面。")

    # --- 4 控制循环 ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_title_bar(s, "感知-规划-行动闭环是智能体控制的核心", prs)
    add_figure_slide(
        s, assets["loop"],
        "Perceive → Reason → Plan → Act → Observe → Reflect",
        "改绘自 Guo et al. POMDP 智能体形式化",
        prs, img_width=Inches(5.0), top=Inches(1.55),
    )
    add_takeaway_strip(s, "与生成式 AI 的本质差异：可观测反馈并修正计划", prs)
    add_footer(s, footer(4), prs)
    add_speaker_notes(s, "控制循环是理解所有后续范式（ReAct、Reflexion）的基础；强调 closed-loop 而非 open-loop 生成。")

    # --- 5 架构模块 ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_title_bar(s, "四模块架构支撑 Profile-Memory-Planning-Action 闭环", prs)
    add_two_column_bullets(
        s,
        "四大核心模块",
        ["Profile：角色与约束", "Memory：短期上下文 + 长期检索", "Planning：目标分解与反思", "Action：工具调用与环境交互"],
        "与生成式 AI 对比",
        ["生成式：prompt → 单次输出", "智能体：多步循环 + 环境反馈", "评估维度：延迟、成本、可靠性、安全", "形式化：POMDP 控制问题 (Guo et al.)"],
    )
    add_footer(s, footer(5), prs)
    add_speaker_notes(s, "四模块是工程实现视角；与上一页三轴能力分类互补——一个偏概念，一个偏系统架构。")

    # --- 6 基础范式 ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_title_bar(s, "ReAct 奠定了思考与行动交织的基础范式", prs)
    add_table_slide(
        s,
        ["范式", "年份", "核心思想", "主要局限"],
        [
            ["ReAct", "2022", "Thought 与 Action 交替", "错误传播"],
            ["Toolformer", "2023", "自监督工具学习", "训练开销大"],
            ["Reflexion", "2023", "语言化自我反思", "推理成本上升"],
            ["AutoGPT", "2023", "自主目标循环", "稳定性不足"],
        ],
    )
    add_footer(s, footer(6), prs)
    add_speaker_notes(s, "范式演进体现同一控制循环的不同实现策略；ReAct 影响最深远，后续工作多在其上叠加记忆、反思或工具学习。")

    # --- 7 演进时间线 ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_title_bar(s, "2022–2025 年技术栈从范式创新走向工程编排", prs)
    add_figure_slide(
        s, assets["timeline"],
        "ReAct → Toolformer → AutoGPT → LangGraph → MCP → GAIA",
        "整理自综述文献时间线",
        prs, img_width=Inches(10.8), top=Inches(1.4),
    )
    add_footer(s, footer(7), prs)
    add_speaker_notes(s, "时间线显示：早期是算法范式突破，2024 起框架与协议（LangGraph、MCP）成为生态关键节点。")

    # --- 8 框架生态 ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_title_bar(s, "LangGraph 等框架推动多智能体编排走向生产", prs)
    add_table_slide(
        s,
        ["框架", "范式", "多智能体", "生产就绪度"],
        [
            ["LangGraph", "状态图编排", "是", "高"],
            ["AutoGen", "多智能体对话", "是", "中"],
            ["CrewAI", "角色分工团队", "是", "中"],
            ["MetaGPT", "软件公司模拟", "是", "中"],
            ["CAMEL", "角色扮演协作", "是", "新兴"],
        ],
        top=Inches(1.4), font_size=11,
    )
    add_footer(s, footer(8), prs)
    add_speaker_notes(s, "框架选择影响拓扑、状态管理与工具接入方式；LangGraph 因显式状态机与可观测性在工业场景更受青睐。")

    # --- 9 框架成熟度 ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_title_bar(s, "框架选择对性能的影响不亚于基础模型", prs)
    add_figure_slide(
        s, assets["framework"],
        "同一模型在不同框架下任务成功率差异显著",
        "改绘自 MASEval, 2026",
        prs, img_width=Inches(9.2), top=Inches(1.35),
    )
    add_takeaway_strip(s, "评估应把框架拓扑作为一等变量，而非仅比较模型", prs)
    add_footer(s, footer(9), prs)
    add_speaker_notes(s, "MASEval 核心论点：现有 benchmark 过于 model-centric；实际部署中框架与拓扑同等重要。")

    # --- 10 协议 ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_title_bar(s, "MCP 标准化工具接入，亦引入新型攻击面", prs)
    add_bullets(s, [
        "MCP：统一 LLM 工具与资源发现接口",
        "ACP / A2A：跨智能体消息与协作协议",
        "收益：降低集成摩擦，模块化工具生态",
        "风险：prompt injection、越权调用、数据外泄",
    ], top=Inches(1.4))
    add_footer(s, footer(10), prs)
    add_speaker_notes(s, "协议层是 2024–2025 新热点；MCP 已被主流 IDE 与 agent 平台采纳，但安全模型仍不成熟。")

    # --- 11 评估基准 ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_title_bar(s, "现有基准偏重模型，忽视框架拓扑效应", prs)
    add_table_slide(
        s,
        ["基准", "领域", "核心指标"],
        [
            ["AgentBench", "OS/DB/Web 等8环境", "成功率"],
            ["GAIA", "通用 AI 助手", "分级准确率"],
            ["SWE-bench", "软件工程", "Issue 解决率"],
            ["τ-bench", "客服智能体", "任务完成率"],
            ["Agent-SafetyBench", "智能体安全", "伤害规避分"],
        ],
        top=Inches(1.35), font_size=11,
    )
    add_takeaway_strip(s, "缺口：缺少系统级、多维（安全/成本/协作）统一基准", prs)
    add_footer(s, footer(11), prs)
    add_speaker_notes(s, "列举主流 benchmark 及其侧重点；指出 MASEval 发现的 model-centric 评估盲区。")

    # --- 12 应用场景 ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_title_bar(s, "软件工程与科研发现是首批规模化落地场景", prs)
    add_two_column_bullets(
        s,
        "行业领域",
        ["软件工程 (SWE-Agent)", "科学发现与文献综合", "医疗与生物医学", "金融与企业自动化", "教育个性化辅导"],
        "典型能力",
        ["自主 GitHub Issue 修复", "多步 Web 深度调研", "带工具调用的客服", "代码生成与 DevOps", "材料科学假设生成"],
    )
    add_footer(s, footer(12), prs)
    add_speaker_notes(s, "应用侧已出现可演示的成功案例，但多数仍属 pilot；高 stakes 领域对可靠性与可验证性要求更高。")

    # --- 13 挑战（可靠性/安全） ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_title_bar(s, "幻觉式工具调用与 prompt 注入仍是部署瓶颈", prs)
    add_bullets(s, [
        "可靠性：工具幻觉、无限循环、级联错误",
        "安全性：注入攻击、越权操作、数据外泄",
        "成本：每步推理增加延迟与 token 开销",
        "可复现性：采样随机性、API 变更、动态 Web 环境",
    ], top=Inches(1.4))
    add_footer(s, footer(13), prs)
    add_speaker_notes(s, "将挑战分为技术可靠性、安全、经济与 reproducibility；这些是阻碍 production 的核心障碍。")

    # --- 14 挑战（评估/伦理）+ 未来 ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_title_bar(s, "系统级评估与协议安全模型是下一阶段重点", prs)
    add_bullets(s, [
        "评估：框架拓扑、安全、成本、协作多维统一",
        "算法：RL 动态工具选择与长程规划",
        "协议：形式化安全模型与标准化互操作",
        "融合：神经-符号推理用于高风险可验证域",
        "协作：人机协同中自主性与问责的平衡",
    ], top=Inches(1.35), size=15)
    add_footer(s, footer(14), prs)
    add_speaker_notes(s, "未来方向呼应前文缺口：评估体系、协议安全、可验证推理、human-in-the-loop 治理。")

    # --- 15 总结 ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, NAVY)
    title = s.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(11.9), Inches(0.9))
    title.text_frame.text = "智能体 AI 正走向工程化，可靠部署仍待突破"
    title.text_frame.paragraphs[0].font.size = Pt(30)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = WHITE
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    takeaways = s.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(3.6))
    tf = takeaways.text_frame
    tf.word_wrap = True
    for i, pt in enumerate([
        "智能体 = LLM 推理 + 工具使用 + 多步交互",
        "ReAct 奠定 think-act-observe 基础循环",
        "框架选择对性能影响不亚于模型选择",
        "生产落地需安全、成本与可靠性工程",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "▸  " + pt
        p.font.size = Pt(17)
        p.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        p.space_after = Pt(12)

    thanks = s.shapes.add_textbox(Inches(0.7), Inches(5.9), Inches(11.9), Inches(0.7))
    thanks.text_frame.text = "谢谢 · 欢迎提问"
    thanks.text_frame.paragraphs[0].font.size = Pt(26)
    thanks.text_frame.paragraphs[0].font.color.rgb = ACCENT
    thanks.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_footer(s, footer(15), prs)
    add_speaker_notes(s, "收束五条 takeaway；强调 field maturing but dependable deployment remains open——与 Nature-style bounded implication 一致。")

    prs.save(OUTPUT)
    audit = audit_presentation(OUTPUT)
    write_qa_report(audit, total)

    print(f"Presentation saved: {OUTPUT}")
    print(f"QA report: {QA_REPORT}")
    print(f"Assets: {ASSETS_DIR}")
    for line in audit:
        print(f"  {line}")


if __name__ == "__main__":
    build_presentation()
